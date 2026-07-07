#!/usr/bin/env python3
"""選択スライドを1つのPPTXに合体する。
   使い方: python3 plus_merge.py <decks_dir> <out_path>
   標準入力に JSON: [{"deck":"ファイル名.pptx","slide":0}, ...]（選択順）"""
import sys, json, copy, io, os
from pptx import Presentation
from pptx.oxml.ns import qn

def del_all_slides(prs):
    # sldIdLst を空にし、presentation が持つ「スライド」relを全て落とす。
    # ※ 過去の編集で sldIdLst から外れたが本体が残った孤立スライドも
    #    reltype で拾って落とす（これを残すとマージ時にパーツ名が衝突する）。
    part = prs.part
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        lst.remove(sid)
    for rId, rel in list(part.rels.items()):
        if rel.reltype.endswith("/slide"):
            part.drop_rel(rId)

def blank_layout(prs):
    best, fewest = None, 99
    for lay in prs.slide_layouts:
        n = len(lay.placeholders)
        if n < fewest:
            fewest, best = n, lay
    return best

def copy_slide(dest, src_slide, layout):
    ns = dest.slides.add_slide(layout)
    for ph in list(ns.placeholders):
        ph._element.getparent().remove(ph._element)
    for sh in src_slide.shapes:
        ns.shapes._spTree.append(copy.deepcopy(sh._element))
    for blip in ns._element.iter(qn('a:blip')):
        rid = blip.get(qn('r:embed'))
        if not rid:
            continue
        try:
            src_part = src_slide.part.related_part(rid)
        except KeyError:
            continue
        _, new_rid = ns.part.get_or_add_image_part(io.BytesIO(src_part.blob))
        blip.set(qn('r:embed'), new_rid)

def main():
    decks_dir, out_path = sys.argv[1], sys.argv[2]
    sel = json.load(sys.stdin)
    if not sel:
        sys.exit("empty selection")
    # ベース＝最初の選択デッキ（同一テンプレなのでどれでも可）
    base = os.path.join(decks_dir, sel[0]["deck"])
    dest = Presentation(base)
    del_all_slides(dest)
    # 一度保存して読み直し、孤立パーツ（元スライド/ノート等）を一掃する。
    # これをしないと新規スライドのパーツ名が旧パーツと衝突し、PowerPointが
    # 「修復」ダイアログを出す原因になる（zip内に同名パーツが二重に入る）。
    _buf = io.BytesIO()
    dest.save(_buf)
    _buf.seek(0)
    dest = Presentation(_buf)
    lay = blank_layout(dest)
    cache = {}
    for it in sel:
        path = os.path.join(decks_dir, it["deck"])
        prs = cache.get(path) or Presentation(path)
        cache[path] = prs
        copy_slide(dest, prs.slides[it["slide"]], lay)
    dest.save(out_path)
    print(json.dumps({"ok": True, "slides": len(sel)}))

if __name__ == "__main__":
    main()
